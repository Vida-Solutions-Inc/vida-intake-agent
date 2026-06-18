"""Turn a file on disk into something the model can read: text or an image block.

Every optional dependency is imported lazily and guarded, so the package still
imports and runs with only the core deps installed — a file we can't extract
simply comes back as a note, and the agent routes it to review.
"""

from __future__ import annotations

import base64
from dataclasses import dataclass
from pathlib import Path

TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".log", ".json", ".xml",
    ".html", ".htm", ".rtf", ".eml", ".vtt", ".srt", ".tex",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif"}
_IMAGE_MEDIA = {
    ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
    ".gif": "image/gif", ".webp": "image/webp",
}
# Anthropic vision works best with the long edge <= ~1568px and payload < ~5MB.
_MAX_IMG_EDGE = 1568
_MAX_IMG_BYTES = 4_500_000


@dataclass
class Extracted:
    kind: str                 # "text" | "image" | "empty" | "binary" | "error"
    text: str = ""
    media_type: str = ""      # for images
    data_b64: str = ""        # for images
    note: str = ""            # human/agent-readable status

    @property
    def is_image(self) -> bool:
        return self.kind == "image"


def extract(path: Path, max_bytes: int = 4_000_000) -> Extracted:
    suffix = path.suffix.lower()
    try:
        if suffix in IMAGE_EXTENSIONS:
            return _extract_image(path)
        if suffix == ".pdf":
            return _extract_pdf(path, max_bytes)
        if suffix == ".docx":
            return _extract_docx(path, max_bytes)
        if suffix == ".xlsx":
            return _extract_xlsx(path, max_bytes)
        if suffix == ".pptx":
            return _extract_pptx(path, max_bytes)
        if suffix in TEXT_EXTENSIONS or _looks_like_text(path):
            return _extract_text(path, max_bytes)
    except Exception as e:  # never let extraction crash the pipeline
        return Extracted("error", note=f"extraction failed: {type(e).__name__}: {e}")

    size = _safe_size(path)
    return Extracted(
        "binary",
        note=f"Binary or unsupported file ({suffix or 'no extension'}, {size} bytes). "
        f"Decide from the filename and folder context, or route to review.",
    )


# ----------------------------------------------------------------- text formats
def _extract_text(path: Path, max_bytes: int) -> Extracted:
    raw = path.read_bytes()[: max_bytes]
    text = raw.decode("utf-8", errors="replace")
    if not text.strip():
        return Extracted("empty", note="File is empty.")
    return Extracted("text", text=text)


def _looks_like_text(path: Path) -> bool:
    try:
        chunk = path.read_bytes()[:2048]
    except OSError:
        return False
    if not chunk:
        return False
    if b"\x00" in chunk:
        return False
    nonprintable = sum(1 for b in chunk if b < 9 or (13 < b < 32))
    return nonprintable / len(chunk) < 0.05


# ----------------------------------------------------------------- PDF
def _extract_pdf(path: Path, max_bytes: int) -> Extracted:
    try:
        from pypdf import PdfReader
    except ImportError:
        return Extracted("binary", note="PDF support not installed (pip install 'intake-agent[extract]').")
    reader = PdfReader(str(path))
    chunks: list[str] = []
    total = 0
    for page in reader.pages:
        t = (page.extract_text() or "").strip()
        if t:
            chunks.append(t)
            total += len(t)
        if total >= max_bytes:
            chunks.append("... (truncated)")
            break
    text = "\n\n".join(chunks).strip()
    if not text:
        return Extracted(
            "empty",
            note=f"PDF has no extractable text ({len(reader.pages)} pages) — likely "
            f"a scan. Decide from the filename, or route to review.",
        )
    return Extracted("text", text=text)


# ----------------------------------------------------------------- DOCX
def _extract_docx(path: Path, max_bytes: int) -> Extracted:
    try:
        import docx
    except ImportError:
        return Extracted("binary", note="DOCX support not installed (pip install 'intake-agent[extract]').")
    document = docx.Document(str(path))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    text = "\n".join(parts).strip()[:max_bytes]
    if not text:
        return Extracted("empty", note="DOCX has no text content.")
    return Extracted("text", text=text)


# ----------------------------------------------------------------- XLSX
def _extract_xlsx(path: Path, max_bytes: int) -> Extracted:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return Extracted("binary", note="XLSX support not installed (pip install 'intake-agent[extract]').")
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    total = 0
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                line = " | ".join(cells)
                parts.append(line)
                total += len(line)
            if total >= max_bytes:
                parts.append("... (truncated)")
                break
        if total >= max_bytes:
            break
    wb.close()
    text = "\n".join(parts).strip()
    if not text:
        return Extracted("empty", note="Spreadsheet has no values.")
    return Extracted("text", text=text)


# ----------------------------------------------------------------- PPTX
def _extract_pptx(path: Path, max_bytes: int) -> Extracted:
    try:
        from pptx import Presentation
    except ImportError:
        return Extracted("binary", note="PPTX support not installed (pip install 'intake-agent[extract]').")
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    line = "".join(r.text for r in p.runs).strip()
                    if line:
                        parts.append(line)
    text = "\n".join(parts).strip()[:max_bytes]
    if not text:
        return Extracted("empty", note="Presentation has no text.")
    return Extracted("text", text=text)


# ----------------------------------------------------------------- images
def _extract_image(path: Path) -> Extracted:
    suffix = path.suffix.lower()
    media = _IMAGE_MEDIA.get(suffix, "image/png")
    # Try to downscale/normalise with Pillow; fall back to raw bytes.
    try:
        from io import BytesIO

        from PIL import Image

        img = Image.open(path)
        img = img.convert("RGB") if img.mode not in ("RGB", "L") else img
        long_edge = max(img.size)
        if long_edge > _MAX_IMG_EDGE:
            scale = _MAX_IMG_EDGE / long_edge
            img = img.resize((int(img.width * scale), int(img.height * scale)))
        buf = BytesIO()
        img.save(buf, format="JPEG", quality=85)
        data = buf.getvalue()
        media = "image/jpeg"
    except Exception:
        data = path.read_bytes()
        if len(data) > _MAX_IMG_BYTES:
            return Extracted(
                "binary",
                note=f"Image too large to send ({len(data)} bytes) and Pillow not "
                f"available to resize. Route from filename or to review.",
            )
    return Extracted("image", media_type=media, data_b64=base64.b64encode(data).decode("ascii"))


def _safe_size(path: Path) -> int:
    try:
        return path.stat().st_size
    except OSError:
        return -1
